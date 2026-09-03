#!/usr/bin/env python3
"""Locally inspect PDF or PPTX presentations without using external services."""

from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import uuid
from pathlib import Path
from typing import Any


DEFAULT_OUTPUT_ROOT = Path("07 Sources/Slides/Processed")
SUPPORTED_SUFFIXES = {".pdf", ".pptx"}


class InspectionError(RuntimeError):
    """A user-facing inspection failure."""


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Extract local technical metadata and per-page/per-slide text from "
            "a PDF or PPTX. The source is never moved or modified."
        )
    )
    parser.add_argument("source", type=Path, help="Path to a PDF or PPTX source")
    parser.add_argument(
        "--output-root",
        type=Path,
        default=DEFAULT_OUTPUT_ROOT,
        help=f"Root for technical results (default: {DEFAULT_OUTPUT_ROOT})",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Replace only manifest.json and extracted-text.md in an existing result",
    )
    return parser.parse_args(argv)


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source_file:
        for chunk in iter(lambda: source_file.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def safe_result_name(stem: str) -> str:
    name = re.sub(r"[\x00-\x1f\x7f]", "", stem)
    name = re.sub(r"[/\\:]+", " - ", name)
    name = re.sub(r"\s+", " ", name).strip(" .")
    if not name:
        name = "presentation"
    return name[:180]


def json_value(value: Any) -> Any:
    if isinstance(value, (dt.date, dt.datetime)):
        return value.isoformat()
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return str(value)


def clean_metadata(metadata: dict[str, Any]) -> dict[str, Any]:
    return {
        str(key): json_value(value)
        for key, value in metadata.items()
        if value not in (None, "")
    }


def inspect_pdf_with_pymupdf(source: Path) -> tuple[dict[str, Any], str]:
    try:
        import pymupdf  # type: ignore
    except ImportError as exc:
        raise InspectionError("PyMuPDF is not installed") from exc

    try:
        document = pymupdf.open(source)
    except Exception as exc:
        raise InspectionError(f"PyMuPDF could not open {source.name}: {exc}") from exc

    try:
        sections: list[str] = []
        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            sections.append(
                f"# Page {page_number}\n\n{text or '[No extractable text]'}"
            )
        manifest = {
            "source": source.name,
            "format": "pdf",
            "pages": document.page_count,
            "backend": "PyMuPDF",
            "metadata": clean_metadata(document.metadata or {}),
        }
        return manifest, "\n\n".join(sections) + "\n"
    finally:
        document.close()


def parse_pdfinfo(output: str) -> dict[str, str]:
    metadata: dict[str, str] = {}
    for line in output.splitlines():
        key, separator, value = line.partition(":")
        if separator and key.strip():
            metadata[key.strip()] = value.strip()
    return metadata


def inspect_pdf_with_poppler(source: Path) -> tuple[dict[str, Any], str]:
    pdfinfo = shutil.which("pdfinfo")
    pdftotext = shutil.which("pdftotext")
    if not pdfinfo or not pdftotext:
        raise InspectionError(
            "PDF inspection requires PyMuPDF or both local pdfinfo and pdftotext"
        )

    info_process = subprocess.run(
        [pdfinfo, str(source)],
        check=False,
        capture_output=True,
        text=True,
    )
    if info_process.returncode != 0:
        raise InspectionError(
            f"pdfinfo could not inspect {source.name}: {info_process.stderr.strip()}"
        )

    text_process = subprocess.run(
        [pdftotext, "-layout", str(source), "-"],
        check=False,
        capture_output=True,
        text=True,
    )
    if text_process.returncode != 0:
        raise InspectionError(
            f"pdftotext could not inspect {source.name}: {text_process.stderr.strip()}"
        )

    metadata = parse_pdfinfo(info_process.stdout)
    try:
        page_count = int(metadata.get("Pages", "0"))
    except ValueError:
        page_count = 0

    pages = text_process.stdout.split("\f")
    if pages and not pages[-1].strip():
        pages.pop()
    if page_count > len(pages):
        pages.extend([""] * (page_count - len(pages)))
    if page_count == 0:
        page_count = len(pages)

    sections = [
        f"# Page {number}\n\n{page.strip() or '[No extractable text]'}"
        for number, page in enumerate(pages, start=1)
    ]
    manifest = {
        "source": source.name,
        "format": "pdf",
        "pages": page_count,
        "backend": "Poppler pdfinfo/pdftotext",
        "metadata": clean_metadata(metadata),
    }
    return manifest, "\n\n".join(sections) + "\n"


def inspect_pdf(source: Path) -> tuple[dict[str, Any], str]:
    try:
        return inspect_pdf_with_pymupdf(source)
    except InspectionError as pymupdf_error:
        try:
            return inspect_pdf_with_poppler(source)
        except InspectionError as poppler_error:
            raise InspectionError(
                f"PDF inspection unavailable. {pymupdf_error}. {poppler_error}."
            ) from poppler_error


def shape_texts(shape: Any) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = getattr(shape, "text", "").strip()
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_values = [cell.text.strip() for cell in row.cells]
            row_text = " | ".join(value for value in row_values if value)
            if row_text:
                texts.append(row_text)
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is not None:
        for child in child_shapes:
            texts.extend(shape_texts(child))
    return texts


def unique_texts(values: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        normalized = value.strip()
        if normalized and normalized not in seen:
            seen.add(normalized)
            result.append(normalized)
    return result


def inspect_pptx(source: Path) -> tuple[dict[str, Any], str]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise InspectionError(
            "PPTX inspection requires python-pptx. Install scripts/slides/requirements.txt "
            "inside a local virtual environment."
        ) from exc

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        raise InspectionError(f"python-pptx could not open {source.name}: {exc}") from exc

    core = presentation.core_properties
    metadata = clean_metadata(
        {
            "title": core.title,
            "author": core.author,
            "subject": core.subject,
            "keywords": core.keywords,
            "category": core.category,
            "comments": core.comments,
            "created": core.created,
            "modified": core.modified,
            "last_modified_by": core.last_modified_by,
            "revision": core.revision,
        }
    )

    sections: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""

        body_parts: list[str] = []
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            body_parts.extend(shape_texts(shape))
        body_parts = unique_texts(body_parts)

        notes = ""
        if getattr(slide, "has_notes_slide", False):
            notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
            if notes_frame is not None:
                notes = notes_frame.text.strip()

        lines = [
            f"# Slide {slide_number}",
            "",
            f"Title: {title or '[No title]'}",
            "",
            "## Text",
            "",
            "\n\n".join(body_parts) if body_parts else "[No extractable text]",
        ]
        if notes:
            lines.extend(["", "## Notes", "", notes])
        sections.append("\n".join(lines))

    manifest = {
        "source": source.name,
        "format": "pptx",
        "slides": len(presentation.slides),
        "backend": "python-pptx",
        "metadata": metadata,
    }
    return manifest, "\n\n".join(sections) + "\n"


def atomic_write(path: Path, content: str) -> None:
    temporary = path.with_name(f".{path.name}.{uuid.uuid4().hex}.tmp")
    try:
        temporary.write_text(content, encoding="utf-8")
        os.replace(temporary, path)
    finally:
        if temporary.exists():
            temporary.unlink()


def write_results(
    output_directory: Path,
    manifest: dict[str, Any],
    extracted_text: str,
    overwrite: bool,
) -> None:
    manifest_path = output_directory / "manifest.json"
    text_path = output_directory / "extracted-text.md"
    existing = [path.name for path in (manifest_path, text_path) if path.exists()]
    if existing and not overwrite:
        names = ", ".join(existing)
        raise InspectionError(
            f"Result already exists in {output_directory}: {names}. "
            "Use --overwrite to replace only these generated files."
        )

    output_directory.mkdir(parents=True, exist_ok=True)
    atomic_write(
        manifest_path,
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
    )
    atomic_write(text_path, extracted_text)


def inspect(source: Path) -> tuple[dict[str, Any], str]:
    suffix = source.suffix.lower()
    if suffix == ".pdf":
        return inspect_pdf(source)
    if suffix == ".pptx":
        return inspect_pptx(source)
    raise InspectionError(
        f"Unsupported format {suffix or '[no extension]'}. Expected PDF or PPTX."
    )


def main(argv: list[str] | None = None) -> int:
    args = parse_args(argv)
    source = args.source.expanduser()
    if not source.exists():
        print(f"error: source does not exist: {source}", file=sys.stderr)
        return 2
    if not source.is_file():
        print(f"error: source is not a file: {source}", file=sys.stderr)
        return 2
    if source.suffix.lower() not in SUPPORTED_SUFFIXES:
        print(
            f"error: unsupported format {source.suffix or '[no extension]'}; "
            "expected .pdf or .pptx",
            file=sys.stderr,
        )
        return 2

    try:
        source = source.resolve(strict=True)
        manifest, extracted_text = inspect(source)
        manifest["sha256"] = file_sha256(source)
        manifest["generated_at"] = dt.datetime.now(dt.timezone.utc).isoformat()
        output_directory = (
            args.output_root.expanduser() / safe_result_name(source.stem)
        )
        write_results(output_directory, manifest, extracted_text, args.overwrite)
    except (InspectionError, OSError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 3

    item_count = manifest.get("pages", manifest.get("slides", 0))
    print(
        f"Inspected {source.name}: {item_count} pages/slides via "
        f"{manifest['backend']}. Result: {output_directory}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
